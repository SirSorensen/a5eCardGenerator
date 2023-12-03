import pptx
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx import Presentation

def printFont(font):
    print("font.name = ", font.name)
    print("font.size.pt = ", font.size.pt, "pt")
    try:
        print("language_id = ", font.language_id)
    except KeyError:
        font.language_id = MSO_LANGUAGE_ID.ENGLISH_US
        print("language_id = ", font.language_id)
    if font.color.type == MSO_COLOR_TYPE.RGB:
        print("font.color.rgb = ", font.color.rgb)
    elif font.color.type == MSO_COLOR_TYPE.SCHEME:
        print("font.color.theme_color = ", font.color.theme_color)
    print("font.fill.type = ", font.fill.type)

    values = " -> "

    if font.bold: values += "bold, "
    if font.italic: values += "italic, "
    if font.underline: values += "underline"

    if len(values) > 4: print(values)

    print()

def main():

    prs = Presentation("cardCreator-test.pptx")


    prs.slides[0].shapes[0].name = 'Background'
    prs.slides[0].shapes[1].name = 'Back Image'
    prs.slides[0].shapes[2].name = 'Back Title'
    prs.slides[0].shapes[3].name = 'Back Subtitle'
    prs.slides[0].shapes[4].name = 'Front Subtitle'
    prs.slides[0].shapes[5].name = 'Front Description'
    prs.slides[0].shapes[6].name = 'Front Title'
    prs.slides[0].shapes[7].name = 'Back Icon 1'

    for shape in prs.slides[0].shapes:
        print("")
        print(prs.slides[0].shapes.index(shape), shape.name, type(shape), shape.width, shape.height)
        if type(shape) is pptx.shapes.picture.Picture:
            print("Picture!")
        elif type(shape) is pptx.shapes.autoshape.Shape:
            print("Autoshape!")
            #print(shape.text)

        if shape.has_text_frame:
            textFrame = shape.text_frame
            print("Number of paragraphs = ", len(textFrame.paragraphs))
            for par in textFrame.paragraphs:
                print(par.text)
                print()

                print("Number of runs = ", len(par.runs))
                for run in par.runs:
                    print(run.text)
                    printFont(run.font)
        else:  
            print("No text_frame :(")




    prs.save('Outputs\\output-test.pptx')
            