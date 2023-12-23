import pptx
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx import Presentation
from pptx.text.text import TextFrame, Font
from pptx.parts.image import Image

class PowerPointInspector:

    def printFont(font):
        print("font.name = ", font.name)
        print("font.size.pt = ", font.size.pt, "pt")
        try:
            print("language_id = ", font.language_id)
        except KeyError:
            font.language_id = MSO_LANGUAGE_ID.ENGLISH_US
            print("language_id = ", font.language_id)
        if font.color.type == MSO_COLOR_TYPE.RGB:
            print("font.color.rgb = ", font.color.rgb)
        elif font.color.type == MSO_COLOR_TYPE.SCHEME:
            print("font.color.theme_color = ", font.color.theme_color)
        print("font.fill.type = ", font.fill.type)

        values = " -> "

        if font.bold: values += "bold, "
        if font.italic: values += "italic, "
        if font.underline: values += "underline"

        if len(values) > 4: print(values)

        print()

    def get_fonts(text_frame : TextFrame) -> list[Font]:
        fonts : list[Font] = []
        pars = text_frame.paragraphs

        for par in pars:
            for run in par.runs:
                if run.font not in fonts:
                    fonts.append(run.font)
        
        return fonts
    
    def get_font_dict(slide : pptx.slide.Slide) -> dict[str, list[Font]]:
        font_dict : dict[str, Font] = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                font_dict[shape.name] = PowerPointInspector.get_fonts(text_frame)

        return font_dict

    def main():
        prs = Presentation("cardCreator-test.pptx")

        template_shapes = prs.slides[0].shapes

        template_shapes[0].name = 'Background'
        template_shapes[1].name = 'Back Image'
        template_shapes[2].name = 'Back Title'
        template_shapes[3].name = 'Back Subtitle'
        template_shapes[4].name = 'Front Subtitle'
        template_shapes[5].name = 'Front Description'
        template_shapes[6].name = 'Front Title'
        template_shapes[7].name = 'Back Icon 1'

        for shape in template_shapes:
            print("")
            print(template_shapes.index(shape), shape.name, type(shape), shape.width, shape.height)
            if type(shape) is pptx.shapes.picture.Picture:
                print("Picture!")
            elif type(shape) is pptx.shapes.autoshape.Shape:
                print("Autoshape!")
                #print(shape.text)

            if shape.has_text_frame:
                textFrame = shape.text_frame
                print("Number of paragraphs = ", len(textFrame.paragraphs))
                for par in textFrame.paragraphs:
                    print(par.text)
                    print()

                    print("Number of runs = ", len(par.runs))
                    runs = len(par.runs)
                    for i in range(runs):
                        run = par.runs[i]
                        print(run.text)
                        PowerPointInspector.printFont(run.font)
                        font_name = run.font.name
                        if i > 1:
                            font_name += f"_{i}"
            else:  
                print("No text_frame :(")

        prs.save('cardCreator-test.pptx')
                