from PowerPoint_Editing.power_point_inspector import PowerPointInspector
from data_forge.data_structures.context_contents import Paragraphs
from data_forge.data_structures.card import Card

import copy

import pptx
from pptx import Presentation
from pptx.text.text import TextFrame
from lxml import etree
from pptx.enum.text import PP_ALIGN

from data_forge.file_handlers.file_handler import FileHandler



class PowerPointModifier:
    def __init__(self):
        self.pptx = Presentation(FileHandler.gen_slide_output_name())
        self.fonts = PowerPointInspector.get_font_dict(self.pptx.slides[0])
        self.current_slide = 0

    def _get_blank_slide_layout(self):
        layout_items_count = [len(layout.placeholders) for layout in self.pptx.slide_layouts]
        min_items = min(layout_items_count)
        blank_layout_id = layout_items_count.index(min_items)
        return self.pptx.slide_layouts[blank_layout_id]

    def duplicate_slide(self, index):
        """Duplicate the slide with the given index in self.

        Adds slide to the end of the presentation"""
        template = self.pptx.slides[index]
        blank_slide_layout = PowerPointModifier._get_blank_slide_layout(self)
        copied_slide = self.pptx.slides.add_slide(blank_slide_layout)

        total_shapes = len(template.shapes)
        for i in range(total_shapes):
            shape = template.shapes[i]
            if type(shape) is pptx.shapes.picture.Picture:
                match shape.name:
                    case 'Background':
                        filename = 'Resources/Images/Background.png'
                    case 'Back Image':
                        filename = 'Resources/Images/Image.png'
                    case 'Back Icon 1':
                        filename = 'Resources/Images/Icon.png'
                    case _:
                        raise ValueError(f"Unknown picture name: {shape.name}")
                    
                copied_slide.shapes.add_picture(filename, shape.left, shape.top, shape.width, shape.height)
                copied_slide.shapes[i].name = shape.name
            else:
                element = shape.element
                newel = copy.deepcopy(element)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            
        return copied_slide
    
    def insertCard(self, card : Card):
        self.duplicate_slide(0)
        self.current_slide += 1
        self.insertFront(card)
        self.insertBack(card)
    
    def insertFront(self, card : Card):
        self.insertTitle(card)
        self.insertSubtitle(card)
        self.insertDescription(card)
    
    def insertBack(self, card : Card):
        self.insertTitle(card, back=True)
        self.insertSubtitle(card, back=True)
        # self.insertIcon(card, back=True)
    
    def insertTitle(self, card : Card, back=False):
        if back:
            shape = self.pptx.slides[self.current_slide].shapes[2]
            self.insert_text(shape, card.title.upper())
        else:
            shape = self.pptx.slides[self.current_slide].shapes[6]
            self.insert_text(shape, card.title.capitalize())
        
    
    def insertSubtitle(self, card : Card, back=False):
        if back:
            shape = self.pptx.slides[self.current_slide].shapes[4]
        else:
            shape = self.pptx.slides[self.current_slide].shapes[3]
        
        self.insert_text(shape, card.subtitle)
    
    def insertDescription(self, card : Card):
        shape = self.pptx.slides[self.current_slide].shapes[5]
        if not isinstance(card.description, Paragraphs):
            raise ValueError(f"Description for {str(type(card))} is not a Paragraphs object, but a {str(type(card.description))} object")
        self.insert_body_text(shape, card.description)
    
    def insertIcon(self, card : Card):
        shape = self.pptx.slides[self.current_slide].shapes[1]
        shape.image = card.icon
    
    def insertImage(self, card : Card):
        shape = self.pptx.slides[self.current_slide].shapes[7]
        shape.image = card.image
    
    def save(self):
        self.pptx.save(FileHandler.gen_slide_output_directory())
    
    def __ready_text_frame(self, shape):
        if not shape.has_text_frame:
            raise ValueError("Shape does not have text frame")
        
        text_frame = shape.text_frame
        fonts = self.fonts[shape.name]
        
        text_frame.clear()  # remove any existing paragraphs, leaving one empty one
        
        par = text_frame.paragraphs[0]

        return fonts, par
    
    def __new_paragraph(self, text_frame : TextFrame):
        if len(text_frame.paragraphs) == 0 or text_frame.paragraphs[0].text != "":
            par = text_frame.add_paragraph()
            par.bullet = False
            par._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
            par.level = 0
            return par
       
        return text_frame.paragraphs[0]

    def insert_text(self, shape, paragraph_str):
        fonts, par = self.__ready_text_frame(shape)

        run = par.add_run()
        run.text = paragraph_str
        PowerPointModifier.set_font(run, fonts[0])
    
    def insert_body_text(self, shape, paragraphs : Paragraphs):        
        fonts, _ = self.__ready_text_frame(shape)
        paragraph_list = paragraphs.to_list()

        for paragraph in paragraph_list:
            par = self.__new_paragraph(shape.text_frame)
            for text in paragraph.text_list:
                run = par.add_run()
                font_int = 0 if text.is_title else 1
                PowerPointModifier.set_font(run, fonts[font_int])
                run.text = str(text)
         

    def set_font(run, font):
        run.font.name = font.name
        run.font.size = font.size
        run.font.bold = font.bold
        run.font.italic = font.italic
        run.font.underline = font.underline
        
        if font.color.type is not None:
            run.font.color.rgb = font.color.rgb